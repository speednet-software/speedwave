"""Integration tests for the office worker's library-driven Python support-scripts.

Each script is invoked the way the TypeScript layer invokes it — `python3 <script> <args...>`,
single JSON object on stdout, non-zero exit on failure — and the produced file is checked for
the right magic bytes and (where cheap) re-opened with the library to confirm it is valid.

These need the worker's Python deps (python-docx, openpyxl, python-pptx, pypdf, matplotlib).
Run with `make test-mcp-office-py` (which builds a venv from `requirements.txt`), or inside the
office image. When the deps are absent the module is skipped, not failed.
"""

from __future__ import annotations

import json
import subprocess
import sys
import zipfile
from pathlib import Path

import pytest

SCRIPTS_DIR = Path(__file__).resolve().parent

# Skip the whole module unless every library a script uses is importable.
for _mod in ("docx", "openpyxl", "pptx", "pypdf", "matplotlib"):
    pytest.importorskip(_mod, reason=f"office Python dependency '{_mod}' not installed")


def _matplotlib_renders() -> bool:
    """Probe whether matplotlib can actually render a figure on this interpreter.

    matplotlib < 3.11 hits a `copy.deepcopy` RecursionError on Python 3.14; the worker image runs
    an older Python where this works, but a local test interpreter may not. Tests that need a real
    render (`renderChart`, image embedding, PDF fixtures) skip when this returns False.
    """
    import io

    import matplotlib

    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    try:
        fig, ax = plt.subplots(figsize=(1, 1))
        ax.plot([0, 1], [0, 1])
        buf = io.BytesIO()
        fig.savefig(buf)
        plt.close(fig)
        return True
    except RecursionError:
        return False


MATPLOTLIB_OK = _matplotlib_renders()
needs_matplotlib = pytest.mark.skipif(
    not MATPLOTLIB_OK, reason="matplotlib cannot render on this interpreter (too-new Python for the pinned matplotlib)"
)


def _weasyprint_loads() -> bool:
    """WeasyPrint dlopens pango/harfbuzz/fontconfig through cffi, raising OSError, not ImportError."""
    import contextlib
    import io

    try:
        # The dlopen failure path prints a multi-line install hint to stdout before raising.
        with contextlib.redirect_stdout(io.StringIO()):
            import weasyprint  # noqa: F401
    except (ImportError, OSError):
        return False
    return True


WEASYPRINT_OK = _weasyprint_loads()
needs_weasyprint = pytest.mark.skipif(
    not WEASYPRINT_OK, reason="weasyprint cannot load its native libraries (pango/harfbuzz/fontconfig)"
)


def _make_pdf(path: Path, pages: int) -> None:
    """Write a `pages`-page blank PDF via pypdf (no matplotlib — works on any interpreter)."""
    from pypdf import PdfWriter

    writer = PdfWriter()
    for _ in range(pages):
        writer.add_blank_page(width=200, height=200)
    with open(path, "wb") as fh:
        writer.write(fh)


def _make_form_pdf(path: Path, field_name: str) -> None:
    """Write a one-page PDF with a single AcroForm text field named `field_name`."""
    from pypdf import PdfWriter
    from pypdf.generic import (
        ArrayObject,
        BooleanObject,
        DictionaryObject,
        NameObject,
        NumberObject,
        TextStringObject,
    )

    writer = PdfWriter()
    page = writer.add_blank_page(width=200, height=200)
    field = DictionaryObject(
        {
            NameObject("/FT"): NameObject("/Tx"),
            NameObject("/T"): TextStringObject(field_name),
            NameObject("/Subtype"): NameObject("/Widget"),
            NameObject("/Rect"): ArrayObject(
                [NumberObject(10), NumberObject(10), NumberObject(100), NumberObject(30)]
            ),
        }
    )
    field_ref = writer._add_object(field)
    page[NameObject("/Annots")] = ArrayObject([field_ref])
    writer._root_object[NameObject("/AcroForm")] = DictionaryObject(
        {
            NameObject("/Fields"): ArrayObject([field_ref]),
            NameObject("/NeedAppearances"): BooleanObject(True),
        }
    )
    with open(path, "wb") as fh:
        writer.write(fh)


def _make_broken_widget_form_pdf(path: Path, field_name: str) -> None:
    """Write a one-page PDF whose single AcroForm widget is missing `/Rect`: pypdf's
    `update_page_form_field_values` raises `KeyError: '/Rect'` writing to it (a genuine failure,
    not the "no fields on this page" case `_fillform`'s per-page except is meant to absorb)."""
    from pypdf import PdfWriter
    from pypdf.generic import ArrayObject, BooleanObject, DictionaryObject, NameObject, TextStringObject

    writer = PdfWriter()
    page = writer.add_blank_page(width=200, height=200)
    field = DictionaryObject(
        {
            NameObject("/FT"): NameObject("/Tx"),
            NameObject("/T"): TextStringObject(field_name),
            NameObject("/Subtype"): NameObject("/Widget"),
        }
    )
    field_ref = writer._add_object(field)
    page[NameObject("/Annots")] = ArrayObject([field_ref])
    writer._root_object[NameObject("/AcroForm")] = DictionaryObject(
        {
            NameObject("/Fields"): ArrayObject([field_ref]),
            NameObject("/NeedAppearances"): BooleanObject(True),
        }
    )
    with open(path, "wb") as fh:
        writer.write(fh)


def _make_nested_form_pdf(path: Path, parent_name: str, child_name: str, *, broken: bool = False) -> None:
    """Write a one-page PDF with a hierarchical AcroForm field: a parent `/T` node with a single
    child text-widget kid. `get_fields()` keys this by the dotted qualified name
    (`f"{parent_name}.{child_name}"`); the widget's own bare `/T` is `child_name`. When `broken` is
    True the widget is missing `/Rect`, so writing to it raises `KeyError: '/Rect'`."""
    from pypdf import PdfWriter
    from pypdf.generic import ArrayObject, BooleanObject, DictionaryObject, NameObject, NumberObject, TextStringObject

    writer = PdfWriter()
    page = writer.add_blank_page(width=200, height=200)
    parent_field = DictionaryObject({NameObject("/T"): TextStringObject(parent_name), NameObject("/Kids"): ArrayObject()})
    parent_ref = writer._add_object(parent_field)
    child: dict = {
        NameObject("/FT"): NameObject("/Tx"),
        NameObject("/T"): TextStringObject(child_name),
        NameObject("/Subtype"): NameObject("/Widget"),
        NameObject("/Parent"): parent_ref,
    }
    if not broken:
        child[NameObject("/Rect")] = ArrayObject([NumberObject(10), NumberObject(10), NumberObject(100), NumberObject(30)])
    child_ref = writer._add_object(DictionaryObject(child))
    parent_field[NameObject("/Kids")] = ArrayObject([child_ref])
    page[NameObject("/Annots")] = ArrayObject([child_ref])
    writer._root_object[NameObject("/AcroForm")] = DictionaryObject(
        {
            NameObject("/Fields"): ArrayObject([parent_ref]),
            NameObject("/NeedAppearances"): BooleanObject(True),
        }
    )
    with open(path, "wb") as fh:
        writer.write(fh)


def run_script(name: str, *args: str) -> dict:
    """Run `scripts/<name>` with the current interpreter; assert exit 0; return the parsed stdout JSON."""
    proc = subprocess.run(
        [sys.executable, str(SCRIPTS_DIR / name), *args],
        capture_output=True,
        text=True,
    )
    assert proc.returncode == 0, f"{name} failed (exit {proc.returncode}): {proc.stderr or proc.stdout}"
    out = json.loads(proc.stdout)
    assert out.get("ok") is True, f"{name} did not report ok: {out}"
    return out


def run_script_expect_fail(name: str, *args: str) -> dict:
    """Run `scripts/<name>`; assert it exits non-zero and reports `{ok: false}`; return the parsed stdout dict."""
    proc = subprocess.run(
        [sys.executable, str(SCRIPTS_DIR / name), *args],
        capture_output=True,
        text=True,
    )
    assert proc.returncode != 0, f"{name} unexpectedly succeeded: {proc.stdout}"
    if not proc.stdout.strip():
        return {}
    result = json.loads(proc.stdout)
    assert result.get("ok") is False
    return result


def is_zip(path: Path) -> bool:
    """True if `path` is a ZIP container (docx/xlsx/pptx are OOXML ZIPs)."""
    return zipfile.is_zipfile(path)


def is_pdf(path: Path) -> bool:
    """True if `path` starts with the PDF magic bytes."""
    return path.read_bytes()[:5] == b"%PDF-"


# ── docx_build.py ────────────────────────────────────────────────────────────


def test_docx_create_all_element_types(tmp_path: Path) -> None:
    out = tmp_path / "doc.docx"
    spec = {
        "elements": [
            {"type": "heading", "level": 1, "text": "Title"},
            {"type": "paragraph", "text": "Body", "bold": True},
            {"type": "table", "header": ["A", "B"], "rows": [["1", "2"], ["3", "4"]]},
            {"type": "pagebreak"},
            {"type": "paragraph", "text": "After break"},
        ]
    }
    res = run_script("docx_build.py", "create", str(out), json.dumps(spec))
    assert res["path"] == str(out)
    assert is_zip(out)
    from docx import Document

    doc = Document(str(out))
    texts = [p.text for p in doc.paragraphs]
    assert "Title" in texts and "Body" in texts and "After break" in texts


def _png_1x1() -> bytes:
    """A minimal valid 1×1 transparent RGBA PNG (correct chunk lengths/CRCs) — no matplotlib/Pillow needed."""
    import struct
    import zlib

    def chunk(typ: bytes, data: bytes) -> bytes:
        return struct.pack(">I", len(data)) + typ + data + struct.pack(">I", zlib.crc32(typ + data) & 0xFFFFFFFF)

    sig = b"\x89PNG\r\n\x1a\n"
    ihdr = chunk(b"IHDR", struct.pack(">IIBBBBB", 1, 1, 8, 6, 0, 0, 0))
    idat = chunk(b"IDAT", zlib.compress(b"\x00\x00\x00\x00\x00"))
    iend = chunk(b"IEND", b"")
    return sig + ihdr + idat + iend


def test_docx_create_with_image(tmp_path: Path) -> None:
    img = tmp_path / "dot.png"
    img.write_bytes(_png_1x1())
    out = tmp_path / "with-img.docx"
    spec = {"elements": [{"type": "image", "path": str(img)}]}
    run_script("docx_build.py", "create", str(out), json.dumps(spec))
    assert is_zip(out)


def test_docx_edit_ops(tmp_path: Path) -> None:
    src = tmp_path / "src.docx"
    run_script(
        "docx_build.py",
        "create",
        str(src),
        json.dumps({"elements": [{"type": "paragraph", "text": "hello world"}, {"type": "paragraph", "text": "keep"}, {"type": "paragraph", "text": "drop"}]}),
    )
    out = tmp_path / "edited.docx"
    ops = [
        {"op": "replace_text", "find": "hello", "replace": "goodbye"},
        {"op": "append", "element": {"type": "paragraph", "text": "appended"}},
        {"op": "delete_paragraph", "index": 2},
    ]
    run_script("docx_build.py", "edit", str(src), str(out), json.dumps(ops))
    from docx import Document

    texts = [p.text for p in Document(str(out)).paragraphs]
    assert any("goodbye world" in t for t in texts)
    assert "appended" in texts
    assert "drop" not in texts


def test_docx_errors(tmp_path: Path) -> None:
    out = tmp_path / "x.docx"
    run_script_expect_fail("docx_build.py", "create", str(out), json.dumps({"elements": [{"type": "video"}]}))
    run_script_expect_fail("docx_build.py", "create", str(out), "not json")
    src = tmp_path / "s.docx"
    run_script("docx_build.py", "create", str(src), json.dumps({"elements": []}))
    run_script_expect_fail("docx_build.py", "edit", str(src), str(out), json.dumps([{"op": "delete_paragraph", "index": 99}]))


def test_docx_replace_text_zero_matches_is_a_reported_failure(tmp_path: Path) -> None:
    src = tmp_path / "src.docx"
    run_script("docx_build.py", "create", str(src), json.dumps({"elements": [{"type": "paragraph", "text": "hello world"}]}))
    out = tmp_path / "edited.docx"
    result = run_script_expect_fail(
        "docx_build.py", "edit", str(src), str(out), json.dumps([{"op": "replace_text", "find": "not-present", "replace": "x"}])
    )
    assert "not-present" in result["error"]
    assert "was not found" in result["error"]
    assert not out.exists()


def test_docx_replace_text_duplicate_op_in_same_batch_is_idempotent(tmp_path: Path) -> None:
    """A second replace_text op with a `find` already consumed by an earlier op in the same batch
    is a no-op success (the text is already gone), not a batch-aborting zero-match failure."""
    src = tmp_path / "src.docx"
    run_script(
        "docx_build.py",
        "create",
        str(src),
        json.dumps({"elements": [{"type": "paragraph", "text": "{{name}}"}]}),
    )
    out = tmp_path / "edited.docx"
    ops = [
        {"op": "replace_text", "find": "{{name}}", "replace": "Acme"},
        {"op": "replace_text", "find": "{{name}}", "replace": "Beta"},
    ]
    run_script("docx_build.py", "edit", str(src), str(out), json.dumps(ops))
    from docx import Document

    texts = [p.text for p in Document(str(out)).paragraphs]
    assert texts == ["Acme"]


def test_docx_replace_text_zero_matches_when_present_only_in_a_table_cell_is_a_reported_failure(tmp_path: Path) -> None:
    """`find` present only inside a table cell (not a paragraph) must still count as a match, not a false zero-match failure."""
    src = tmp_path / "src.docx"
    run_script(
        "docx_build.py",
        "create",
        str(src),
        json.dumps({"elements": [{"type": "table", "header": ["A"], "rows": [["only-in-cell"]]}]}),
    )
    out = tmp_path / "edited.docx"
    run_script(
        "docx_build.py",
        "edit",
        str(src),
        str(out),
        json.dumps([{"op": "replace_text", "find": "only-in-cell", "replace": "replaced"}]),
    )
    from docx import Document

    doc = Document(str(out))
    cell_texts = [cell.text for table in doc.tables for row in table.rows for cell in row.cells]
    assert any("replaced" in t for t in cell_texts)


# ── xlsx_build.py ────────────────────────────────────────────────────────────


def test_xlsx_create_with_chart_and_freeze(tmp_path: Path) -> None:
    out = tmp_path / "wb.xlsx"
    spec = {
        "sheets": [
            {
                "name": "Data",
                "rows": [["Month", "Sales"], ["Jan", 10], ["Feb", 20], ["Mar", 30]],
                "freeze": "A2",
                "charts": [
                    {"type": "bar", "title": "Sales", "dataRange": "Data!B1:B4", "categoriesRange": "Data!A2:A4", "anchor": "D2"}
                ],
            },
            {"name": "Empty", "rows": []},
        ]
    }
    run_script("xlsx_build.py", "create", str(out), json.dumps(spec))
    assert is_zip(out)
    from openpyxl import load_workbook

    wb = load_workbook(str(out))
    assert wb.sheetnames == ["Data", "Empty"]
    assert wb["Data"]["B2"].value == 10
    assert wb["Data"].freeze_panes == "A2"


def test_xlsx_edit_ops(tmp_path: Path) -> None:
    src = tmp_path / "src.xlsx"
    run_script("xlsx_build.py", "create", str(src), json.dumps({"sheets": [{"name": "S", "rows": [[1, 2]]}]}))
    out = tmp_path / "edited.xlsx"
    ops = [
        {"op": "set_cell", "sheet": "S", "cell": "A3", "value": 7},
        {"op": "set_formula", "sheet": "S", "cell": "B3", "formula": "SUM(A1:A3)"},
        {"op": "add_sheet", "name": "Extra"},
        {"op": "add_chart", "sheet": "S", "chart": {"type": "line", "dataRange": "S!A1:A3", "anchor": "D1"}},
    ]
    run_script("xlsx_build.py", "edit", str(src), str(out), json.dumps(ops))
    from openpyxl import load_workbook

    wb = load_workbook(str(out))
    assert "Extra" in wb.sheetnames
    assert wb["S"]["A3"].value == 7
    assert str(wb["S"]["B3"].value).startswith("=SUM")


def test_xlsx_errors(tmp_path: Path) -> None:
    out = tmp_path / "x.xlsx"
    run_script_expect_fail("xlsx_build.py", "create", str(out), json.dumps({"sheets": []}))
    run_script_expect_fail("xlsx_build.py", "create", str(out), json.dumps({"sheets": [{"name": "S", "rows": [], "charts": [{"type": "donut", "dataRange": "S!A1:A2", "anchor": "B1"}]}]}))


def test_xlsx_unknown_sheet_name_is_a_teaching_error(tmp_path: Path) -> None:
    src = tmp_path / "src.xlsx"
    run_script("xlsx_build.py", "create", str(src), json.dumps({"sheets": [{"name": "Q1", "rows": [[1, 2]]}]}))
    out = tmp_path / "edited.xlsx"
    result = run_script_expect_fail(
        "xlsx_build.py", "edit", str(src), str(out), json.dumps([{"op": "set_cell", "sheet": "Q1 Data", "cell": "A1", "value": 5}])
    )
    assert "Q1 Data" in result["error"]
    assert "['Q1']" in result["error"]


def test_xlsx_unknown_sheet_name_is_a_teaching_error_via_set_formula(tmp_path: Path) -> None:
    src = tmp_path / "src.xlsx"
    run_script("xlsx_build.py", "create", str(src), json.dumps({"sheets": [{"name": "Q1", "rows": [[1, 2]]}]}))
    out = tmp_path / "edited.xlsx"
    result = run_script_expect_fail(
        "xlsx_build.py", "edit", str(src), str(out), json.dumps([{"op": "set_formula", "sheet": "Q1 Data", "cell": "A1", "formula": "=SUM(A1:A2)"}])
    )
    assert "Q1 Data" in result["error"]
    assert "['Q1']" in result["error"]


def test_xlsx_unknown_sheet_name_is_a_teaching_error_via_add_chart(tmp_path: Path) -> None:
    src = tmp_path / "src.xlsx"
    run_script("xlsx_build.py", "create", str(src), json.dumps({"sheets": [{"name": "Q1", "rows": [[1, 2]]}]}))
    out = tmp_path / "edited.xlsx"
    result = run_script_expect_fail(
        "xlsx_build.py",
        "edit",
        str(src),
        str(out),
        json.dumps([{"op": "add_chart", "sheet": "Q1 Data", "chart": {"type": "bar", "dataRange": "Q1!B1:B2", "anchor": "D2"}}]),
    )
    assert "Q1 Data" in result["error"]
    assert "['Q1']" in result["error"]


# ── pptx_build.py ────────────────────────────────────────────────────────────


def test_pptx_create_bullets_and_chart(tmp_path: Path) -> None:
    out = tmp_path / "deck.pptx"
    spec = {
        "slides": [
            {"title": "Agenda", "bullets": ["First", "Second", "Third"]},
            {"title": "Numbers", "chart": {"type": "column", "categories": ["Q1", "Q2"], "series": [{"name": "Revenue", "values": [100.0, 150.0]}], "title": "Revenue"}},
        ]
    }
    run_script("pptx_build.py", "create", str(out), json.dumps(spec))
    assert is_zip(out)
    from pptx import Presentation

    prs = Presentation(str(out))
    assert len(prs.slides) == 2


def test_pptx_edit_ops(tmp_path: Path) -> None:
    src = tmp_path / "src.pptx"
    run_script("pptx_build.py", "create", str(src), json.dumps({"slides": [{"title": "One"}, {"title": "Two"}]}))
    out = tmp_path / "edited.pptx"
    ops = [
        {"op": "add_slide", "slide": {"title": "Three", "bullets": ["x"]}},
        {"op": "set_title", "index": 0, "text": "Renamed"},
        {"op": "delete_slide", "index": 1},
    ]
    run_script("pptx_build.py", "edit", str(src), str(out), json.dumps(ops))
    from pptx import Presentation

    prs = Presentation(str(out))
    titles = [s.shapes.title.text if s.shapes.title else None for s in prs.slides]
    assert "Renamed" in titles
    assert "Three" in titles
    assert "Two" not in titles


def test_pptx_errors(tmp_path: Path) -> None:
    out = tmp_path / "x.pptx"
    run_script_expect_fail("pptx_build.py", "create", str(out), json.dumps({"slides": []}))
    src = tmp_path / "s.pptx"
    run_script("pptx_build.py", "create", str(src), json.dumps({"slides": [{"title": "a"}]}))
    run_script_expect_fail("pptx_build.py", "edit", str(src), str(out), json.dumps([{"op": "delete_slide", "index": 99}]))


# ── render_chart.py ──────────────────────────────────────────────────────────


@needs_matplotlib
@pytest.mark.parametrize("ctype", ["bar", "line", "pie", "scatter", "area"])
def test_render_chart_png(tmp_path: Path, ctype: str) -> None:
    out = tmp_path / f"{ctype}.png"
    spec = {
        "type": ctype,
        "title": ctype.title(),
        "data": {"labels": ["A", "B", "C"], "series": [{"name": "s1", "values": [1.0, 2.0, 3.0]}]},
    }
    run_script("render_chart.py", str(out), json.dumps(spec))
    data = out.read_bytes()
    assert data[:8] == b"\x89PNG\r\n\x1a\n"
    assert len(data) > 100  # a complete (atomically-written) PNG, not a truncated one
    # No `.tmp-<uuid>` sibling left behind by atomic_save.
    assert not list(tmp_path.glob(f"{ctype}.png.tmp-*"))


@needs_matplotlib
def test_render_chart_svg_and_multiseries(tmp_path: Path) -> None:
    out = tmp_path / "chart.svg"
    spec = {
        "type": "line",
        "format": "svg",
        "data": {"labels": ["A", "B"], "series": [{"name": "x", "values": [1, 2]}, {"name": "y", "values": [3, 4]}]},
    }
    run_script("render_chart.py", str(out), json.dumps(spec))
    assert b"<svg" in out.read_bytes()[:512]


def test_render_chart_errors(tmp_path: Path) -> None:
    # Validation happens before matplotlib is touched, so these run on any interpreter.
    out = tmp_path / "x.png"
    run_script_expect_fail("render_chart.py", str(out), json.dumps({"type": "donut", "data": {"labels": ["A"], "series": [{"name": "s", "values": [1]}]}}))
    run_script_expect_fail("render_chart.py", str(out), json.dumps({"type": "bar", "data": {"labels": ["A", "B"], "series": [{"name": "s", "values": [1]}]}}))


# ── pdf_ops.py ───────────────────────────────────────────────────────────────


def test_pdf_metadata_merge_split_rotate(tmp_path: Path) -> None:
    a, b = tmp_path / "a.pdf", tmp_path / "b.pdf"
    _make_pdf(a, 3)
    _make_pdf(b, 2)

    meta = run_script("pdf_ops.py", "metadata", str(a))["metadata"]
    assert meta["pages"] == 3
    assert meta["encrypted"] is False
    # A blank PDF with no title/author/creator set reports null (not "" or a missing key);
    # pypdf stamps its own /Producer when writing, so that field alone is non-null here.
    for key in ("title", "author", "creator"):
        assert meta[key] is None

    merged = tmp_path / "merged.pdf"
    res = run_script("pdf_ops.py", "merge", str(merged), str(a), str(b))
    assert res["pages"] == 5
    assert is_pdf(merged)

    part = tmp_path / "part.pdf"
    res = run_script("pdf_ops.py", "split", str(a), str(part), "1", "2")
    assert res["pages"] == 2
    assert is_pdf(part)

    rotated = tmp_path / "rot.pdf"
    res = run_script("pdf_ops.py", "rotate", str(a), str(rotated), "90", "1,3")
    assert res["rotated"] == [1, 3]
    assert is_pdf(rotated)


def test_pdf_watermark_and_fillform(tmp_path: Path) -> None:
    doc, wm = tmp_path / "doc.pdf", tmp_path / "wm.pdf"
    _make_pdf(doc, 2)
    _make_pdf(wm, 1)
    out = tmp_path / "watermarked.pdf"
    res = run_script("pdf_ops.py", "watermark", str(doc), str(wm), str(out))
    assert res["pages"] == 2
    assert is_pdf(out)

    # fillform on a PDF with no AcroForm fields is a no-op per page but still succeeds.
    filled = tmp_path / "filled.pdf"
    res = run_script("pdf_ops.py", "fillform", str(doc), str(filled), "1", json.dumps({"name": "Ada"}))
    assert is_pdf(filled)


def test_pdf_fillform_warns_on_unknown_field_name(tmp_path: Path) -> None:
    form = tmp_path / "form.pdf"
    _make_form_pdf(form, "applicant_name")
    out = tmp_path / "filled.pdf"
    res = run_script(
        "pdf_ops.py",
        "fillform",
        str(form),
        str(out),
        "0",
        json.dumps({"applicant_name": "Ada", "does_not_exist": "x"}),
    )
    assert is_pdf(out)
    assert res["fill_warnings"] == ["unknown field name: 'does_not_exist' — not present in this PDF's AcroForm"]


def test_pdf_fillform_no_warnings_when_all_fields_known(tmp_path: Path) -> None:
    form = tmp_path / "form.pdf"
    _make_form_pdf(form, "applicant_name")
    out = tmp_path / "filled.pdf"
    res = run_script(
        "pdf_ops.py", "fillform", str(form), str(out), "0", json.dumps({"applicant_name": "Ada"})
    )
    assert is_pdf(out)
    assert res["fill_warnings"] is None


def test_pdf_fillform_genuine_write_failure_on_a_matching_page_is_not_silently_absorbed(tmp_path: Path) -> None:
    form = tmp_path / "broken-form.pdf"
    _make_broken_widget_form_pdf(form, "applicant_name")
    out = tmp_path / "filled.pdf"
    result = run_script_expect_fail(
        "pdf_ops.py", "fillform", str(form), str(out), "0", json.dumps({"applicant_name": "Ada"})
    )
    assert "applicant_name" in result["error"]
    assert "KeyError" in result["error"]
    # Must not be misreported as the unrelated "no fields on this page" case.
    assert "no AcroForm fields found" not in result["error"]
    assert not out.exists()


def test_pdf_fillform_hierarchical_field_bare_name_is_not_reported_unknown(tmp_path: Path) -> None:
    """A caller-supplied bare name for a nested (/Parent-chained) field is written correctly by
    pypdf even though get_fields() only exposes the dotted qualified name; it must not be warned
    about as unknown."""
    form = tmp_path / "nested-form.pdf"
    _make_nested_form_pdf(form, "section", "applicant_name")
    out = tmp_path / "filled.pdf"
    res = run_script(
        "pdf_ops.py", "fillform", str(form), str(out), "0", json.dumps({"applicant_name": "Ada"})
    )
    assert is_pdf(out)
    assert res["fill_warnings"] is None

    from pypdf import PdfReader

    reader = PdfReader(str(out))
    fields = reader.get_fields() or {}
    assert fields["section.applicant_name"].get("/V") == "Ada"


def test_pdf_fillform_hierarchical_field_write_failure_is_not_silently_absorbed(tmp_path: Path) -> None:
    """A genuine write failure on a page whose only matching field is nested (name only reachable
    by walking the full /Parent chain) must still surface via fail(), not be swallowed as the
    unrelated 'no fields on this page' case."""
    form = tmp_path / "broken-nested-form.pdf"
    _make_nested_form_pdf(form, "section", "applicant_name", broken=True)
    out = tmp_path / "filled.pdf"
    result = run_script_expect_fail(
        "pdf_ops.py", "fillform", str(form), str(out), "0", json.dumps({"applicant_name": "Ada"})
    )
    assert "applicant_name" in result["error"]
    assert "KeyError" in result["error"]
    assert "no AcroForm fields found" not in result["error"]
    assert not out.exists()


def test_pdf_errors(tmp_path: Path) -> None:
    p = tmp_path / "p.pdf"
    _make_pdf(p, 2)
    out = tmp_path / "x.pdf"
    run_script_expect_fail("pdf_ops.py", "merge", str(out), str(p))  # needs ≥2
    run_script_expect_fail("pdf_ops.py", "split", str(p), str(out), "1", "99")  # out of range
    run_script_expect_fail("pdf_ops.py", "rotate", str(p), str(out), "45", "1")  # bad degrees
    run_script_expect_fail("pdf_ops.py", "bogus")


def test_pdf_merge_bad_input_names_its_position_in_the_batch(tmp_path: Path) -> None:
    good = tmp_path / "good.pdf"
    _make_pdf(good, 1)
    bad = tmp_path / "bad.pdf"
    bad.write_text("not a pdf")
    out = tmp_path / "merged.pdf"
    result = run_script_expect_fail("pdf_ops.py", "merge", str(out), str(good), str(bad))
    assert "item 2 of 2" in result["error"]


def test_pdf_non_pdf_input_is_a_teaching_error(tmp_path: Path) -> None:
    not_a_pdf = tmp_path / "fake.pdf"
    not_a_pdf.write_text("<html><body>this is not a PDF</body></html>")
    result = run_script_expect_fail("pdf_ops.py", "metadata", str(not_a_pdf))
    assert "could not read" in result["error"]
    assert "valid, non-corrupted PDF" in result["error"]


def test_pdf_fillform_non_pdf_input_is_a_teaching_error(tmp_path: Path) -> None:
    not_a_pdf = tmp_path / "fake.pdf"
    not_a_pdf.write_text("<html><body>this is not a PDF</body></html>")
    out = tmp_path / "filled.pdf"
    result = run_script_expect_fail(
        "pdf_ops.py", "fillform", str(not_a_pdf), str(out), "0", json.dumps({"name": "Ada"})
    )
    assert "could not read" in result["error"]
    assert "valid, non-corrupted PDF" in result["error"]


# ── python_docx_extract.py ───────────────────────────────────────────────────


def test_python_docx_extract(tmp_path: Path) -> None:
    src = tmp_path / "src.docx"
    run_script(
        "docx_build.py",
        "create",
        str(src),
        json.dumps({"elements": [{"type": "heading", "level": 2, "text": "Section"}, {"type": "paragraph", "text": "Body text"}, {"type": "table", "header": ["A", "B"], "rows": [["1", "2"]]}]}),
    )
    res = run_script("python_docx_extract.py", str(src))
    md = res["markdown"]
    assert "## Section" in md
    assert "Body text" in md
    assert "| A | B |" in md


def test_python_docx_extract_usage_error(tmp_path: Path) -> None:
    run_script_expect_fail("python_docx_extract.py")  # missing arg


# ── weasyprint_render.py ─────────────────────────────────────────────────────


def test_weasyprint_render_usage_error(tmp_path: Path) -> None:
    # Usage validation runs before weasyprint is imported, so this works on any interpreter.
    run_script_expect_fail("weasyprint_render.py")  # no args
    run_script_expect_fail("weasyprint_render.py", "only-one-arg")


@needs_weasyprint
def test_weasyprint_render_html_to_pdf(tmp_path: Path) -> None:
    src = tmp_path / "page.html"
    src.write_text("<html><head><style>@page{size:A4;margin:18mm}</style></head><body><h1>Hi</h1></body></html>")
    dst = tmp_path / "out.pdf"
    run_script("weasyprint_render.py", str(src), str(dst), f"file://{tmp_path}/")
    assert is_pdf(dst)


@needs_weasyprint
def test_weasyprint_render_rejects_remote_resource(tmp_path: Path) -> None:
    src = tmp_path / "page.html"
    src.write_text('<html><body><img src="https://example.com/x.png"></body></html>')
    dst = tmp_path / "out.pdf"
    # WeasyPrint itself only warns on a fetcher rejection; the script records it and fails.
    run_script_expect_fail("weasyprint_render.py", str(src), str(dst), f"file://{tmp_path}/")
    assert not dst.exists(), "a rejected render must not leave an output PDF"


@needs_weasyprint
def test_weasyprint_render_ignores_presentational_hint_background(tmp_path: Path) -> None:
    src = tmp_path / "page.html"
    src.write_text('<html><body background="http://presentational-hint.invalid/bg.png">hi</body></html>')
    dst = tmp_path / "out.pdf"
    run_script("weasyprint_render.py", str(src), str(dst), f"file://{tmp_path}/")
    assert is_pdf(dst)


@needs_weasyprint
def test_weasyprint_render_ignores_presentational_hint_css_injection(tmp_path: Path) -> None:
    src = tmp_path / "page.html"
    payload = "x);}body{background-image:url(http://presentational-hint.invalid/leak)}"
    src.write_text(f'<html><body background="{payload}">hi</body></html>')
    dst = tmp_path / "out.pdf"
    run_script("weasyprint_render.py", str(src), str(dst), f"file://{tmp_path}/")
    assert is_pdf(dst)
