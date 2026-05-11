"""Create or edit a .pptx from the office worker's slides/ops DSL (python-pptx).

Usage:
  python3 pptx_build.py create <output.pptx> <json-spec>
  python3 pptx_build.py edit   <input.pptx> <output.pptx> <json-ops>

`json-spec` = {"slides": Slide[]}; `json-ops` = Op[]. See ADR-055 for the shapes.
Output: {"ok": true, "path": "<output>"}
"""

from __future__ import annotations

import json

from script_runner import main, ok, fail


_CHART_XL = {
    "column": "COLUMN_CLUSTERED",
    "line": "LINE",
    "pie": "PIE",
    "xy": "XY_SCATTER",
    "bubble": "BUBBLE",
}


def _add_chart(slide, chart_spec: dict) -> None:
    """Add a native chart to ``slide`` from a DSL chart spec."""
    from pptx.chart.data import CategoryChartData, XyChartData, BubbleChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches

    kind = chart_spec.get("type")
    xl_name = _CHART_XL.get(kind)
    if xl_name is None:
        fail(f"unknown chart type: {kind}")
    chart_type = getattr(XL_CHART_TYPE, xl_name)
    x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(5)

    if kind in ("xy", "bubble"):
        # For xy/bubble we treat `categories` as the X values.
        cats = [float(c) for c in chart_spec["categories"]]
        cd = XyChartData() if kind == "xy" else BubbleChartData()
        for ser in chart_spec["series"]:
            s = cd.add_series(str(ser["name"]))
            for cx_val, y_val in zip(cats, ser["values"]):
                if kind == "xy":
                    s.add_data_point(cx_val, float(y_val))
                else:
                    s.add_data_point(cx_val, float(y_val), 1.0)
        graphic_frame = slide.shapes.add_chart(chart_type, x, y, cx, cy, cd)
    else:
        cd = CategoryChartData()
        cd.categories = [str(c) for c in chart_spec["categories"]]
        for ser in chart_spec["series"]:
            cd.add_series(str(ser["name"]), [float(v) for v in ser["values"]])
        graphic_frame = slide.shapes.add_chart(chart_type, x, y, cx, cy, cd)

    if chart_spec.get("title"):
        graphic_frame.chart.has_title = True
        graphic_frame.chart.chart_title.text_frame.text = str(chart_spec["title"])


def _add_slide(prs, slide_spec: dict) -> None:
    """Append a slide built from a DSL slide spec."""
    from pptx.util import Inches

    has_body = bool(slide_spec.get("bullets"))
    layout = prs.slide_layouts[1] if has_body else prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    if slide_spec.get("title") and slide.shapes.title is not None:
        slide.shapes.title.text = str(slide_spec["title"])
    if has_body:
        # Layout 1 has a body placeholder at index 1.
        body = slide.placeholders[1].text_frame
        bullets = list(slide_spec["bullets"])
        body.text = str(bullets[0]) if bullets else ""
        for line in bullets[1:]:
            p = body.add_paragraph()
            p.text = str(line)
    if slide_spec.get("image"):
        slide.shapes.add_picture(str(slide_spec["image"]["path"]), Inches(1), Inches(1.5), height=Inches(5))
    if slide_spec.get("chart"):
        _add_chart(slide, slide_spec["chart"])


def _create(output: str, spec: dict) -> None:
    from pptx import Presentation

    prs = Presentation()
    slides = spec.get("slides", [])
    if not slides:
        fail("spec.slides must be non-empty")
    for slide_spec in slides:
        _add_slide(prs, slide_spec)
    prs.save(output)
    ok(path=output)


def _edit(src: str, output: str, ops: list) -> None:
    from pptx import Presentation

    prs = Presentation(src)
    for op in ops:
        kind = op.get("op")
        if kind == "add_slide":
            _add_slide(prs, op["slide"])
        elif kind == "set_title":
            idx = int(op["index"])
            if idx < 0 or idx >= len(prs.slides):
                fail(f"set_title index out of range: {idx}")
            slide = prs.slides[idx]
            if slide.shapes.title is None:
                fail(f"slide {idx} has no title placeholder")
            slide.shapes.title.text = str(op["text"])
        elif kind == "delete_slide":
            idx = int(op["index"])
            slides = prs.slides
            if idx < 0 or idx >= len(slides):
                fail(f"delete_slide index out of range: {idx}")
            # python-pptx has no public slide-delete; manipulate the XML id list.
            xml_slides = prs.slides._sldIdLst
            slide_ids = list(xml_slides)
            xml_slides.remove(slide_ids[idx])
        else:
            fail(f"unknown op: {kind}")
    prs.save(output)
    ok(path=output)


def _run(argv: list[str]) -> None:
    if not argv:
        fail("usage: pptx_build.py create|edit ...")
    mode = argv[0]
    if mode == "create":
        if len(argv) != 3:
            fail("usage: pptx_build.py create <output.pptx> <json-spec>")
        _create(argv[1], json.loads(argv[2]))
    elif mode == "edit":
        if len(argv) != 4:
            fail("usage: pptx_build.py edit <input.pptx> <output.pptx> <json-ops>")
        _edit(argv[1], argv[2], json.loads(argv[3]))
    else:
        fail(f"unknown mode: {mode}")


if __name__ == "__main__":
    main(_run)
